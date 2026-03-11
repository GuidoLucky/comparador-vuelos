#!/usr/bin/env python3
import sys, json, os, subprocess, tempfile, base64, shutil
from datetime import datetime

def editar_voucher(datos, template_path, out_path):
    from pptx import Presentation
    prs = Presentation(template_path)
    slide = prs.slides[0]
    mapa = {
        'Google Shape;56;p13': datos.get('fecha', datetime.now().strftime('%d/%m/%y')),
        'Google Shape;58;p13': datos.get('servicio', ''),
        'Google Shape;59;p13': datos.get('beneficiario', ''),
        'Google Shape;62;p13': datos.get('fecha_in', ''),
        'Google Shape;63;p13': datos.get('fecha_out', ''),
        'Google Shape;64;p13': datos.get('noches', ''),
        'Google Shape;67;p13': datos.get('hotel', ''),
        'Google Shape;68;p13': datos.get('codigo', ''),
    }
    for shape in slide.shapes:
        if shape.name in mapa and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ''
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = mapa[shape.name]
    prs.save(out_path)

def editar_confirmacion(datos, template_path, out_path):
    import openpyxl
    wb = openpyxl.load_workbook(template_path)
    ws = wb.active
    ws['J4']  = datos.get('fecha', datetime.now().strftime('%d/%m/%Y'))
    ws['C14'] = datos.get('hotel', '')
    ws['D16'] = datos.get('titular', '')
    ws['D18'] = datos.get('pasajeros', '')
    ws['D19'] = datos.get('fecha_in', '')
    ws['G20'] = datos.get('noches', '')
    ws['D21'] = datos.get('fecha_out', '')
    ws['D23'] = datos.get('habitaciones', '-')
    ws['D25'] = datos.get('observaciones', 'N/A')
    ws['D27'] = datos.get('total', '')
    wb.save(out_path)

def convertir_a_pdf(input_path, output_dir):
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf',
         '--outdir', output_dir, input_path],
        capture_output=True, text=True, timeout=60
    )
    if result.returncode != 0:
        raise Exception(f"LibreOffice error: {result.stderr}")
    base = os.path.splitext(os.path.basename(input_path))[0]
    pdf_path = os.path.join(output_dir, base + '.pdf')
    if not os.path.exists(pdf_path):
        pdfs = [f for f in os.listdir(output_dir) if f.endswith('.pdf')]
        if not pdfs:
            raise Exception("No se generó PDF")
        pdf_path = os.path.join(output_dir, pdfs[0])
    return pdf_path

def main():
    if len(sys.argv) < 3:
        print(json.dumps({'ok': False, 'error': 'Uso: generar_docs.py <tipo> <json>'}))
        sys.exit(1)
    tipo = sys.argv[1]
    try:
        datos = json.loads(sys.argv[2])
    except Exception as e:
        print(json.dumps({'ok': False, 'error': f'JSON inválido: {e}'}))
        sys.exit(1)
    base_dir = os.path.dirname(os.path.abspath(__file__))
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            if tipo == 'voucher':
                template = os.path.join(base_dir, 'Editable_Voucher__1_.pptx')
                work_file = os.path.join(tmpdir, 'voucher.pptx')
                shutil.copy2(template, work_file)
                editar_voucher(datos, template, work_file)
            elif tipo == 'confirmacion':
                template = os.path.join(base_dir, 'Confirmación_de_reserva_hoteleria__1_.xlsx')
                if not os.path.exists(template):
                    template = os.path.join(base_dir, 'Confirmacion_de_reserva_hoteleria__1_.xlsx')
                work_file = os.path.join(tmpdir, 'confirmacion.xlsx')
                shutil.copy2(template, work_file)
                editar_confirmacion(datos, template, work_file)
            else:
                print(json.dumps({'ok': False, 'error': f'Tipo desconocido: {tipo}'}))
                sys.exit(1)
            pdf_path = convertir_a_pdf(work_file, tmpdir)
            with open(pdf_path, 'rb') as f:
                pdf_b64 = base64.b64encode(f.read()).decode()
            print(json.dumps({'ok': True, 'pdf': pdf_b64}))
        except Exception as e:
            print(json.dumps({'ok': False, 'error': str(e)}))
            sys.exit(1)

if __name__ == '__main__':
    main()
